import fitz
import chromadb
from chromadb.utils.embedding_functions import SentenceTransformerEmbeddingFunction
from pptx import Presentation
import re
import hashlib
import logging
from pathlib import Path
from functools import lru_cache

from .config import CHROMA_PATH, EMBED_MODEL

logger = logging.getLogger(__name__)

# ChromaDB Setup — lru_cache gives the same run-once semantics as
# st.cache_resource but without coupling the backend to Streamlit
@lru_cache(maxsize=1)
def getClient():
    return chromadb.PersistentClient(path=CHROMA_PATH)

@lru_cache(maxsize=1)
def getEmbedFunction():
    return SentenceTransformerEmbeddingFunction(model_name=EMBED_MODEL)

def get_collection(collectionName: str):
    return getClient().get_or_create_collection(
        name = collectionName,
        embedding_function=getEmbedFunction(),
        metadata={"hnsw:space": "cosine"}
    )

# Text Extraction
def extractFromPDF(path: str):
    pages = []
    document = fitz.open(path)
    for i, page in enumerate(document):
        text = page.get_text()
        text = re.sub(r'\s+', ' ', text) # Clean random spaces
        text = re.sub(r'(?<!\n)\n(?!\n)', ' ', text) # Clean new line that are not a paragraph break
        if text:
            pages.append({"page": i + 1,
                          "text": text})
    return pages

def extractFromPPT(path: str):
    prs = Presentation(path)
    pages = []

    for i, slide in enumerate(prs.slides):
        parts = []

        # Slide title, used as a context prefix for this slide's chunks
        title = ""
        try:
            if slide.shapes.title is not None and slide.shapes.title.text.strip():
                title = slide.shapes.title.text.strip()
        except (AttributeError, KeyError):
            pass

        for shape in slide.shapes:
            # Text frames (titles, body text, text boxes)
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = " ".join(run.text for run in para.runs).strip()
                    if line:
                        parts.append(line)

            # Tables
            if shape.has_table:
                for row in shape.table.rows:
                    cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                    if cells:
                        parts.append(" | ".join(cells))

        text = re.sub(r'\s+', ' ', " ".join(parts)).strip()
        if text:
            pages.append({"page": i + 1, "text": text, "title": title})

    return pages

# Sentence-aware Chunking
def _splitSentences(text: str, maxChars: int) -> list:
    """Split text on sentence boundaries; hard-split any sentence longer than maxChars."""
    parts = re.split(r'(?<=[.!?])\s+', text)
    sentences = []
    for part in parts:
        part = part.strip()
        if not part:
            continue
        # A "sentence" longer than a whole chunk (e.g. a table row or slide
        # with no punctuation) gets hard-split so packing can't stall
        while len(part) > maxChars:
            sentences.append(part[:maxChars])
            part = part[maxChars:].strip()
        if part:
            sentences.append(part)
    return sentences


def textChunk(text: str, source: str, page: int, maxChars: int = 1000,
              overlapChars: int = 150, title: str = ""):
    """Pack whole sentences into ~maxChars chunks with sentence-level overlap.

    Chunks never cut mid-sentence, so each one is a coherent passage for both
    embedding and retrieval. `title` (e.g. a slide title) is prepended to each
    chunk to give the embedding extra context.
    """
    prefix = f"{title}: " if title else ""
    sentences = _splitSentences(text, maxChars)

    chunks = []
    current = []
    length = 0
    idx = 0

    def emit():
        nonlocal idx
        chunk_text = (prefix + " ".join(current)).strip()
        if chunk_text:
            uid = hashlib.md5(f"{source}:{page}:{idx}".encode()).hexdigest()
            chunks.append({
                "uid": uid,
                "text": chunk_text,
                "metadata": {
                    "source": source,
                    "page": page,
                    "index": idx
                }
            })
            idx += 1

    for sentence in sentences:
        if current and length + len(sentence) > maxChars:
            emit()
            # Carry the last sentence(s) over as overlap so answers spanning
            # a chunk boundary are still retrievable in one piece
            kept = []
            kept_len = 0
            for prev in reversed(current):
                if kept_len + len(prev) > overlapChars:
                    break
                kept.insert(0, prev)
                kept_len += len(prev)
            current = kept
            length = kept_len
        current.append(sentence)
        length += len(sentence) + 1

    if current:
        emit()

    return chunks

# Main RAG Pipeline
def fileUpload(filepath: str, collection):
    path = Path(filepath)
    if not path.exists():
        raise FileNotFoundError(f"File not found: {filepath}")

    suffix = path.suffix.lower()

    if suffix == ".pdf":
        pages = extractFromPDF(str(path))
    elif suffix == ".pptx":
        pages = extractFromPPT(str(path))
    elif suffix == ".ppt":
        raise ValueError("Legacy .ppt files are not supported. Save it as .pptx and try again.")
    else:
        raise ValueError("File must be in PDF or PPTX format.")

    if not pages:
        raise ValueError("No text could be extracted from this file.")

    all_chunks = []
    for page in pages:
        all_chunks.extend(textChunk(page["text"], path.name, page["page"],
                                    title=page.get("title", "")))

    # Remove stale chunks from any previous upload of the same file, so a
    # shorter re-upload doesn't leave orphaned chunks in the collection
    try:
        collection.delete(where={"source": path.name})
    except Exception:
        logger.debug("No previous chunks to delete for %s", path.name)

    collection.upsert(
        ids = [chunk["uid"] for chunk in all_chunks],
        documents = [chunk["text"] for chunk in all_chunks],
        metadatas = [chunk["metadata"] for chunk in all_chunks]
    )

    logger.info("%s successfully uploaded to session collection", path.name)