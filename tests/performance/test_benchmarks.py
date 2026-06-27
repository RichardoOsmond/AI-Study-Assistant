"""
Backend function benchmarks using pytest-benchmark.
Tests pure Python logic only — no LLM calls, no network, no Streamlit.
"""
import json
import sys
import pytest
from pathlib import Path
from unittest.mock import MagicMock

ROOT = Path(__file__).parent.parent.parent
sys.path.insert(0, str(ROOT))

from backend.rag import textChunk, extractFromPDF, extractFromPPT
from backend.quiz import _parse_json
from backend.session import loadSessions, createSession
from unittest.mock import patch
from pptx import Presentation
from pptx.util import Inches
import tempfile
import os

SAMPLE_TEXT = """
Artificial intelligence (AI) is intelligence demonstrated by machines, as opposed to
natural intelligence displayed by animals including humans. AI research has been defined
as the field of study of intelligent agents, which refers to any system that perceives
its environment and takes actions that maximize its chance of achieving its goals.
""".strip() * 40  # ~1600 chars repeated to simulate a real document chunk


SAMPLE_QUIZ_JSON = json.dumps([
    {"type": "mcq", "question": "What is AI?", "options": {"A": "...", "B": "...", "C": "...", "D": "..."}, "answer": "A", "explanation": "..."},
    {"type": "short", "question": "Explain machine learning.", "model_answer": "...", "key_points": ["point1", "point2"]},
] * 5)  # 10 questions


# Text Chunking
def test_benchmark_text_chunk_small(benchmark):
    """Benchmark chunking a small ~400 char text."""
    text = SAMPLE_TEXT[:400]
    benchmark(textChunk, text, source="test.pdf", page=1)


def test_benchmark_text_chunk_medium(benchmark):
    """Benchmark chunking a medium ~4000 char text (approx 10 pages)."""
    text = SAMPLE_TEXT[:4000]
    benchmark(textChunk, text, source="test.pdf", page=1)


def test_benchmark_text_chunk_large(benchmark):
    """Benchmark chunking a large ~16000 char text (approx 40 pages)."""
    benchmark(textChunk, SAMPLE_TEXT, source="test.pdf", page=1)


# JSON Parsing
def test_benchmark_parse_json_clean(benchmark):
    """Benchmark parsing a clean JSON array (no markdown fences)."""
    result = benchmark(_parse_json, SAMPLE_QUIZ_JSON)
    assert len(result) == 10


def test_benchmark_parse_json_with_fences(benchmark):
    """Benchmark parsing a JSON array wrapped in markdown fences."""
    fenced = f"```json\n{SAMPLE_QUIZ_JSON}\n```"
    result = benchmark(_parse_json, fenced)
    assert len(result) == 10


# Session Loading
def test_benchmark_load_sessions(benchmark):
    """Benchmark reading and parsing the sessions.json file."""
    result = benchmark(loadSessions)
    assert isinstance(result, dict)


# Chunk Selection
def test_benchmark_chunk_selection(benchmark):
    """Benchmark random chunk selection from a mocked ChromaDB collection."""
    from backend.quiz import chunkSelection

    mock_collection = MagicMock()
    mock_collection.count.return_value = 100
    mock_collection.get.return_value = {
        "documents": [f"Chunk number {i}: " + SAMPLE_TEXT[:200] for i in range(100)]
    }

    result = benchmark(chunkSelection, mock_collection, n=20)
    assert isinstance(result, str)
    assert len(result) > 0


# Session Creation
def test_benchmark_create_session(benchmark, tmp_path):
    """Benchmark creating a new session (UUID generation + JSON read/write)."""
    import backend.session as session_module
    tmp_file = tmp_path / "sessions.json"
    with patch.object(session_module, "sessionsFile", tmp_file):
        result = benchmark(createSession, "Benchmark Session")
    assert result["session_id"] is not None


def test_benchmark_get_sessions_list(benchmark, tmp_path):
    """Benchmark listing sessions with 20 existing sessions."""
    import backend.session as session_module
    tmp_file = tmp_path / "sessions.json"
    with patch.object(session_module, "sessionsFile", tmp_file):
        for i in range(20):
            createSession(f"Session {i}")
        result = benchmark(session_module.getSessionsList)
    assert len(result) == 20


# PPTX Extraction
@pytest.fixture(scope="module")
def sample_pptx_path(tmp_path_factory):
    """Generate a .pptx with 20 slides of varying content for benchmarking."""
    tmp = tmp_path_factory.mktemp("bench_pptx")
    path = tmp / "bench.pptx"
    prs = Presentation()
    blank = prs.slide_layouts[6]

    for i in range(20):
        slide = prs.slides.add_slide(blank)
        tb = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(7), Inches(4))
        tb.text_frame.text = f"Slide {i + 1}: " + SAMPLE_TEXT[:300]

    prs.save(str(path))
    return str(path)


def test_benchmark_extract_pptx_small(benchmark, tmp_path):
    """Benchmark extracting text from a 5-slide .pptx."""
    path = tmp_path / "small.pptx"
    prs = Presentation()
    for i in range(5):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        tb = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(7), Inches(4))
        tb.text_frame.text = f"Slide {i + 1}: " + SAMPLE_TEXT[:300]
    prs.save(str(path))
    result = benchmark(extractFromPPT, str(path))
    assert len(result) == 5


def test_benchmark_extract_pptx_large(benchmark, sample_pptx_path):
    """Benchmark extracting text from a 20-slide .pptx."""
    result = benchmark(extractFromPPT, sample_pptx_path)
    assert len(result) == 20
