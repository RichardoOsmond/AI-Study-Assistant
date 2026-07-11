"""
Tests for extractFromPPT in backend/rag.py.
"""
import sys
import pytest
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches

ROOT = Path(__file__).parent.parent
sys.path.insert(0, str(ROOT))

from backend.rag import extractFromPPT, textChunk, fileUpload


# ── Fixture ───────────────────────────────────────────────────────────────────

@pytest.fixture(scope="module")
def sample_pptx(tmp_path_factory) -> Path:
    """Create a .pptx with known content"""
    tmp = tmp_path_factory.mktemp("pptx")
    path = tmp / "sample.pptx"
    prs = Presentation()
    blank_layout = prs.slide_layouts[6]  # blank layout

    # Slide 1 — title + body text
    slide1 = prs.slides.add_slide(prs.slide_layouts[1])
    slide1.shapes.title.text = "Introduction to AI"
    slide1.placeholders[1].text = "AI stands for Artificial Intelligence."

    # Slide 2 — text box
    slide2 = prs.slides.add_slide(blank_layout)
    txBox = slide2.shapes.add_textbox(Inches(1), Inches(1), Inches(6), Inches(2))
    txBox.text_frame.text = "Machine learning is a subset of AI."

    # Slide 3 — table
    slide3 = prs.slides.add_slide(blank_layout)
    table = slide3.shapes.add_table(2, 2, Inches(1), Inches(1), Inches(5), Inches(2)).table
    table.cell(0, 0).text = "Term"
    table.cell(0, 1).text = "Definition"
    table.cell(1, 0).text = "Neural Network"
    table.cell(1, 1).text = "A model inspired by the human brain."

    # Slide 4 — empty slide (should be skipped)
    prs.slides.add_slide(blank_layout)

    prs.save(str(path))
    return path


@pytest.fixture(scope="module")
def empty_pptx(tmp_path_factory) -> Path:
    """A .pptx with no text on any slide."""
    tmp = tmp_path_factory.mktemp("pptx_empty")
    path = tmp / "empty.pptx"
    prs = Presentation()
    prs.slides.add_slide(prs.slide_layouts[6])
    prs.save(str(path))
    return path


# ── extractFromPPT ────────────────────────────────────────────────────────────

def test_returns_list(sample_pptx):
    result = extractFromPPT(str(sample_pptx))
    assert isinstance(result, list)


def test_skips_empty_slides(sample_pptx):
    """Slide 4 is blank — should not appear in results."""
    result = extractFromPPT(str(sample_pptx))
    assert len(result) == 3


def test_page_numbers_are_correct(sample_pptx):
    result = extractFromPPT(str(sample_pptx))
    assert result[0]["page"] == 1
    assert result[1]["page"] == 2
    assert result[2]["page"] == 3


def test_slide1_contains_title_and_body(sample_pptx):
    result = extractFromPPT(str(sample_pptx))
    text = result[0]["text"]
    assert "Introduction to AI" in text
    assert "Artificial Intelligence" in text


def test_slide2_contains_textbox(sample_pptx):
    result = extractFromPPT(str(sample_pptx))
    assert "Machine learning" in result[1]["text"]


def test_slide3_contains_table(sample_pptx):
    result = extractFromPPT(str(sample_pptx))
    text = result[2]["text"]
    assert "Neural Network" in text
    assert "Definition" in text


def test_empty_pptx_returns_empty_list(empty_pptx):
    result = extractFromPPT(str(empty_pptx))
    assert result == []


def test_each_page_has_required_keys(sample_pptx):
    result = extractFromPPT(str(sample_pptx))
    for page in result:
        assert "page" in page
        assert "text" in page
        assert isinstance(page["text"], str)
        assert len(page["text"]) > 0


# ── Integration: extractFromPPT → textChunk ───────────────────────────────────

def test_chunks_produced_from_pptx(sample_pptx):
    pages = extractFromPPT(str(sample_pptx))
    all_chunks = []
    for page in pages:
        all_chunks.extend(textChunk(page["text"], "sample.pptx", page["page"]))
    assert len(all_chunks) > 0
    for chunk in all_chunks:
        assert "uid" in chunk
        assert "text" in chunk
        assert "metadata" in chunk


# ── Integration: fileUpload with .pptx ───────────────────────────────────────

def test_file_upload_pptx(sample_pptx):
    """fileUpload should accept .pptx and upsert chunks into the collection."""
    from unittest.mock import MagicMock
    mock_collection = MagicMock()
    fileUpload(str(sample_pptx), mock_collection)
    mock_collection.upsert.assert_called_once()
    _, kwargs = mock_collection.upsert.call_args
    assert len(kwargs["ids"]) > 0
    assert len(kwargs["documents"]) == len(kwargs["ids"])


def test_file_upload_unsupported_format(tmp_path):
    """fileUpload should raise ValueError for unsupported file types."""
    from unittest.mock import MagicMock
    bad_file = tmp_path / "notes.txt"
    bad_file.write_text("some text")
    with pytest.raises(ValueError):
        fileUpload(str(bad_file), MagicMock())


def test_file_upload_missing_file():
    """fileUpload should raise FileNotFoundError for non-existent paths."""
    from unittest.mock import MagicMock
    with pytest.raises(FileNotFoundError):
        fileUpload("nonexistent.pptx", MagicMock())


def test_file_upload_empty_pptx_raises(empty_pptx):
    """fileUpload should raise (not silently succeed) when no text is extractable."""
    from unittest.mock import MagicMock
    with pytest.raises(ValueError):
        fileUpload(str(empty_pptx), MagicMock())


def test_file_upload_legacy_ppt_rejected(tmp_path):
    """Legacy .ppt files get a clear error instead of a cryptic parser failure."""
    from unittest.mock import MagicMock
    fake_ppt = tmp_path / "old.ppt"
    fake_ppt.write_bytes(b"\xd0\xcf\x11\xe0 fake legacy ppt")
    with pytest.raises(ValueError, match="pptx"):
        fileUpload(str(fake_ppt), MagicMock())


def test_file_upload_clears_stale_chunks(sample_pptx):
    """Re-uploading a file deletes its previous chunks before upserting."""
    from unittest.mock import MagicMock
    mock_collection = MagicMock()
    fileUpload(str(sample_pptx), mock_collection)
    mock_collection.delete.assert_called_once_with(where={"source": "sample.pptx"})
