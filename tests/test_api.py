"""
FastAPI endpoint tests using TestClient — no server, no browser, no LLM.
AI endpoints are tested for their error paths; happy paths that need Groq
are covered by usage, not CI.
"""
import sys
import pytest
from pathlib import Path
from unittest.mock import MagicMock, patch

ROOT = Path(__file__).parent.parent
sys.path.insert(0, str(ROOT))

from fastapi.testclient import TestClient

import backend.session as session_module
from api.main import app

client = TestClient(app)


@pytest.fixture(autouse=True)
def isolated_sessions(tmp_path):
    """Every test gets its own sessions.json — never touches real data."""
    tmp_file = tmp_path / "sessions.json"
    with patch.object(session_module, "sessionsFile", tmp_file):
        yield


@pytest.fixture
def session_id():
    return client.post("/api/sessions", json={"name": "Test Session"}).json()["session_id"]


# ── Sessions CRUD ─────────────────────────────────────────────────────────────

def test_list_sessions_empty():
    response = client.get("/api/sessions")
    assert response.status_code == 200
    assert response.json() == []

def test_create_session():
    response = client.post("/api/sessions", json={"name": "My Study"})
    assert response.status_code == 201
    body = response.json()
    assert body["display_name"] == "My Study"
    assert body["files"] == []

def test_create_session_empty_name_rejected():
    assert client.post("/api/sessions", json={"name": ""}).status_code == 422

def test_get_session(session_id):
    response = client.get(f"/api/sessions/{session_id}")
    assert response.status_code == 200
    assert response.json()["session_id"] == session_id

def test_get_missing_session_404():
    assert client.get("/api/sessions/nope").status_code == 404

def test_rename_session(session_id):
    response = client.patch(f"/api/sessions/{session_id}", json={"name": "Renamed"})
    assert response.status_code == 200
    assert response.json()["display_name"] == "Renamed"

def test_delete_session(session_id):
    assert client.delete(f"/api/sessions/{session_id}").status_code == 204
    assert client.get(f"/api/sessions/{session_id}").status_code == 404


# ── Files ─────────────────────────────────────────────────────────────────────

def test_upload_rejects_bad_extension(session_id):
    response = client.post(
        f"/api/sessions/{session_id}/files",
        files={"file": ("notes.txt", b"some text", "text/plain")},
    )
    assert response.status_code == 400

def test_upload_pptx(session_id, tmp_path):
    from pptx import Presentation
    from pptx.util import Inches
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tb = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(6), Inches(2))
    tb.text_frame.text = "Neural networks are function approximators. They learn from data."
    pptx_path = tmp_path / "deck.pptx"
    prs.save(str(pptx_path))

    mock_collection = MagicMock()
    with patch("api.main.get_collection", return_value=mock_collection):
        response = client.post(
            f"/api/sessions/{session_id}/files",
            files={"file": ("deck.pptx", pptx_path.read_bytes(),
                            "application/vnd.openxmlformats-officedocument.presentationml.presentation")},
        )

    assert response.status_code == 201
    assert "deck.pptx" in response.json()["files"]
    mock_collection.upsert.assert_called_once()

def test_upload_duplicate_conflict(session_id):
    session_module.addFileSession(session_id, "deck.pptx")
    response = client.post(
        f"/api/sessions/{session_id}/files",
        files={"file": ("deck.pptx", b"x", "application/octet-stream")},
    )
    assert response.status_code == 409

def test_remove_file(session_id):
    session_module.addFileSession(session_id, "deck.pptx")
    mock_collection = MagicMock()
    with patch("api.main.get_collection", return_value=mock_collection):
        response = client.delete(f"/api/sessions/{session_id}/files/deck.pptx")
    assert response.status_code == 200
    assert response.json()["files"] == []
    mock_collection.delete.assert_called_once_with(where={"source": "deck.pptx"})

def test_remove_missing_file_404(session_id):
    assert client.delete(f"/api/sessions/{session_id}/files/nope.pdf").status_code == 404


# ── AI endpoints (non-LLM paths) ──────────────────────────────────────────────

def test_summarize_no_files_400(session_id):
    empty_collection = MagicMock()
    empty_collection.count.return_value = 0
    with patch("api.main.get_collection", return_value=empty_collection):
        response = client.post(f"/api/sessions/{session_id}/summarize",
                               json={"detail_level": "brief"})
    assert response.status_code == 400

def test_summarize_invalid_level(session_id):
    response = client.post(f"/api/sessions/{session_id}/summarize",
                           json={"detail_level": "extreme"})
    assert response.status_code == 400

def test_quiz_generate_no_files_400(session_id):
    empty_collection = MagicMock()
    empty_collection.count.return_value = 0
    with patch("api.main.get_collection", return_value=empty_collection):
        response = client.post(f"/api/sessions/{session_id}/quiz/generate", json={})
    assert response.status_code == 400

def test_quiz_grade_mcq(session_id):
    """MCQ grading needs no LLM — full happy path works through the API."""
    questions = [
        {"type": "mcq", "question": "Q1?", "options": {"A": "1", "B": "2"},
         "answer": "A", "explanation": "because"},
        {"type": "mcq", "question": "Q2?", "options": {"A": "1", "B": "2"},
         "answer": "B", "explanation": "since"},
    ]
    response = client.post(
        f"/api/sessions/{session_id}/quiz/grade",
        json={"questions": questions, "answers": {"0": "A", "1": "A"}},
    )
    assert response.status_code == 200
    body = response.json()
    assert body["score"] == 1
    assert body["total"] == 2
    assert body["results"][0]["correct"] is True
    assert body["results"][1]["correct"] is False

def test_chat_missing_session_404():
    response = client.post("/api/sessions/nope/chat",
                           json={"question": "hi", "history": []})
    assert response.status_code == 404
